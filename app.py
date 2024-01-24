from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import base64

app = Flask(__name__)

def hex_to_rgb(hex_color):
    """ Convert hex color string to RGB. """
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))

@app.route('/create_pptx', methods=['POST'])
def create_pptx():
    data = request.json
    prs = Presentation()

    for slide_data in data.get('slides', []):
        slide_layout = prs.slide_layouts[slide_data['layout']]
        slide = prs.slides.add_slide(slide_layout)

        # Set background color if provided
        if 'background_color' in slide_data:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(slide_data['background_color'])

        # Handling title
        if 'title' in slide_data and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            if 'title_color' in slide_data:
                title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(slide_data['title_color'])
            if 'title_bold' in slide_data:
                title_shape.text_frame.paragraphs[0].font.bold = slide_data['title_bold']
            if 'title_italic' in slide_data:
                title_shape.text_frame.paragraphs[0].font.italic = slide_data['title_italic']

        # Handling content for layouts that have a second placeholder
        if 'content' in slide_data and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            if hasattr(content_placeholder, 'text'):
                content_placeholder.text = slide_data['content']
                if 'content_color' in slide_data:
                    content_placeholder.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(slide_data['content_color'])
                if 'content_bold' in slide_data:
                    content_placeholder.text_frame.paragraphs[0].font.bold = slide_data['content_bold']
                if 'content_italic' in slide_data:
                    content_placeholder.text_frame.paragraphs[0].font.italic = slide_data['content_italic']

        # Adding an image if path is provided and valid
        if 'image_path' in slide_data and os.path.exists(slide_data['image_path']):
            left = Inches(2)
            top = Inches(2)
            width = Inches(4)
            height = Inches(3)
            slide.shapes.add_picture(slide_data['image_path'], left, top, width=width, height=height)

    pptx_file = 'presentation.pptx'
    prs.save(pptx_file)
    with open(pptx_file, "rb") as file:
        encoded_string = base64.b64encode(file.read()).decode()

    return jsonify({"file_data": encoded_string})

if __name__ == '__main__':
    app.run(debug=True)
